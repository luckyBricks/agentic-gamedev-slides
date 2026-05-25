#!/usr/bin/env python3
"""Build editable PPTX from Slidev DOM extraction.

Reads build/extract/slide-NN.json + build/extract/asset-map.json and
emits one editable PowerPoint deck (16:9, 13.333" x 7.5") that mirrors
each slide's text, images, and decorative boxes in slide-local coords.

Run: python3 build/build_pptx.py
"""
from __future__ import annotations

import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# ---- constants -----------------------------------------------------------

EXTRACT_DIR = Path("build/extract")
ASSET_DIR = Path("build/png-assets")
OUT_FILE = Path("build/output/rider-agentic-gamedev.pptx")

SLIDE_W_PX = 1280
SLIDE_H_PX = 720
SLIDE_W_EMU = 12192000  # 13.333"
SLIDE_H_EMU = 6858000   # 7.5"
PX_TO_EMU = SLIDE_W_EMU / SLIDE_W_PX  # 9525
PX_TO_PT = 0.75                       # 96dpi → 72dpi
SCALE = 3                             # asset key uses scaled px

# Generic system stack — we map every Slidev font to Inter explicitly.
FALLBACK_FONT = "Inter"

# ---- helpers -------------------------------------------------------------

def px_emu(v: float) -> int:
    return int(round(v * PX_TO_EMU))


def parse_color(css: str) -> tuple[tuple[int, int, int], float] | None:
    """Return ((r, g, b), alpha 0..1) or None for transparent."""
    if not css:
        return None
    m = re.match(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)(?:\s*,\s*([0-9.]+))?\s*\)", css)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        if a == 0:
            return None
        return (r, g, b), a
    m = re.match(r"#([0-9a-fA-F]{6})", css)
    if m:
        v = m.group(1)
        return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)), 1.0
    if css.startswith("#") and len(css) == 4:
        v = css[1:]
        return (int(v[0] * 2, 16), int(v[1] * 2, 16), int(v[2] * 2, 16)), 1.0
    return None


def blend_over(fg: tuple[tuple[int, int, int], float], bg_rgb: tuple[int, int, int]) -> RGBColor:
    """Alpha-composite fg over solid bg. PPTX shape fill has no alpha channel via the
    python-pptx public API, so we precompute the visible color and use solid fill."""
    (r, g, b), a = fg
    br, bg_, bb = bg_rgb
    out = (
        int(round(r * a + br * (1 - a))),
        int(round(g * a + bg_ * (1 - a))),
        int(round(b * a + bb * (1 - a))),
    )
    return RGBColor(*out)


def parse_radius(css: str) -> int:
    if not css or css == "0px":
        return 0
    m = re.match(r"([0-9.]+)px", css)
    if m:
        return int(round(float(m.group(1))))
    return 0


def parse_line_height(lh: str, font_size_px: float) -> float | None:
    if not lh or lh == "normal":
        return None
    m = re.match(r"([0-9.]+)px", lh)
    if m:
        return float(m.group(1)) / font_size_px
    try:
        return float(lh)
    except (TypeError, ValueError):
        return None


def parse_letter_spacing(ls: str) -> float | None:
    if not ls or ls == "normal":
        return None
    m = re.match(r"(-?[0-9.]+)px", ls)
    if m:
        return float(m.group(1))
    return None


def weight_is_bold(w: str) -> bool:
    try:
        return int(w) >= 600
    except (TypeError, ValueError):
        return w in ("bold", "bolder")


def align_for(css: str) -> PP_ALIGN:
    return {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
        "justify": PP_ALIGN.JUSTIFY,
        "start": PP_ALIGN.LEFT,
        "end": PP_ALIGN.RIGHT,
    }.get((css or "left").lower(), PP_ALIGN.LEFT)


def remove_textbox_inset(tf):
    """Zero out text frame internal padding so positions match Slidev exactly."""
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def asset_key(src: str, w: float, h: float, filter_css: str | None = None) -> str:
    return f"{src}|{int(round(w * SCALE))}|{int(round(h * SCALE))}|{filter_css or ''}"


# ---- builder -------------------------------------------------------------

def build():
    asset_map_path = EXTRACT_DIR / "asset-map.json"
    asset_map: dict[str, str] = json.loads(asset_map_path.read_text())

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)
    blank_layout = prs.slide_layouts[6]  # blank

    slide_files = sorted(EXTRACT_DIR.glob("slide-*.json"))
    print(f"Building {len(slide_files)} slides…")

    for sf in slide_files:
        data = json.loads(sf.read_text())
        slide_n = data.get("slide")
        slide = prs.slides.add_slide(blank_layout)

        # ---- background --------------------------------------------------
        bg = parse_color(data.get("bodyBackground") or data.get("background") or "")
        # Slide-canvas background RGB (used to alpha-blend every translucent overlay).
        slide_bg_rgb: tuple[int, int, int] = (26, 20, 36)
        if bg:
            (r, g, b), _a = bg
            slide_bg_rgb = (r, g, b)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*slide_bg_rgb)

        # ---- decorative boxes -------------------------------------------
        boxes = data.get("boxes", [])
        # Skip the full-slide root box (covers everything) — it's the background already.
        for b in sorted(boxes, key=lambda x: (x.get("zIndex", 0), boxes.index(x))):
            x, y, w, h = b["x"], b["y"], b["w"], b["h"]
            if w >= SLIDE_W_PX - 1 and h >= SLIDE_H_PX - 1:
                continue  # full-slide canvas
            if w < 1 or h < 1:
                continue
            radius = parse_radius(b.get("borderRadius", "")) if b.get("borderRadius") else 0
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
            shp = slide.shapes.add_shape(
                shape_type,
                Emu(px_emu(x)),
                Emu(px_emu(y)),
                Emu(px_emu(w)),
                Emu(px_emu(h)),
            )
            # Adjust rounded-rectangle corner radius — PPTX adjustment is fraction of shorter side.
            if radius > 0:
                try:
                    short = min(w, h)
                    frac = max(0.0, min(0.5, radius / short))
                    shp.adjustments[0] = frac
                except Exception:
                    pass

            # Fill: prefer solid background, else first color of gradient, else nofill.
            # Translucent overlays are alpha-blended against the slide canvas so they
            # render at their *visible* color in PPTX (which lacks fill alpha via API).
            bg_color = parse_color(b.get("bg") or "")
            grad = b.get("gradient")
            if bg_color:
                shp.fill.solid()
                shp.fill.fore_color.rgb = blend_over(bg_color, slide_bg_rgb)
            elif grad:
                m = re.findall(r"#([0-9a-fA-F]{6})|rgba?\(([^)]+)\)", grad)
                first_color = None
                for hex_c, rgb_c in m:
                    if hex_c:
                        first_color = (
                            (int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)),
                            1.0,
                        )
                        break
                    if rgb_c:
                        parts = [p.strip() for p in rgb_c.split(",")]
                        if len(parts) >= 3:
                            alpha = float(parts[3]) if len(parts) >= 4 else 1.0
                            first_color = (
                                (int(parts[0]), int(parts[1]), int(parts[2])),
                                alpha,
                            )
                            break
                if first_color is not None:
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = blend_over(first_color, slide_bg_rgb)
                else:
                    shp.fill.background()
            else:
                shp.fill.background()

            # Border
            bc = parse_color(b.get("borderColor") or "")
            bw = b.get("borderWidth") or 0
            if bc and bw > 0:
                shp.line.color.rgb = blend_over(bc, slide_bg_rgb)
                shp.line.width = Emu(int(round(bw * PX_TO_EMU)))
            else:
                shp.line.fill.background()  # no border

            # No text inside the decorative shape (text is added separately).
            shp.text_frame.text = ""

        # ---- images ------------------------------------------------------
        for im in data.get("images", []):
            key = asset_key(im["src"], im["w"], im["h"], im.get("filter"))
            fname = asset_map.get(key)
            if not fname:
                continue
            fpath = ASSET_DIR / fname
            if not fpath.exists():
                continue
            try:
                slide.shapes.add_picture(
                    str(fpath),
                    Emu(px_emu(im["x"])),
                    Emu(px_emu(im["y"])),
                    width=Emu(px_emu(im["w"])),
                    height=Emu(px_emu(im["h"])),
                )
            except Exception as e:
                print(f"  ! slide {slide_n} image {fname}: {e}")

        # ---- text --------------------------------------------------------
        for t in data.get("texts", []):
            x, y, w, h = t["x"], t["y"], t["w"], t["h"]
            if w < 1 or h < 1:
                continue
            # Skip vertical-writing decorative watermarks — PPTX would wrap each
            # character to its own line in the narrow horizontal box.
            wm = (t.get("writingMode") or "horizontal-tb").lower()
            if wm.startswith("vertical") or wm == "sideways-rl" or wm == "sideways-lr":
                continue
            # Skip near-invisible decorative text (e.g., faint watermarks).
            if (t.get("opacity") is not None) and t["opacity"] < 0.15:
                continue
            fs = t.get("fontSize", 16)
            # Single-line if box height is barely more than one line-height.
            single_line = h <= fs * 1.6 and "\n" not in t.get("text", "")
            # Add a small buffer to width so PPTX's slightly different font metrics
            # don't force an unintended wrap on tight boxes.
            w_render = w + (max(6, fs * 0.6) if single_line else 4)
            tb = slide.shapes.add_textbox(
                Emu(px_emu(x)),
                Emu(px_emu(y)),
                Emu(px_emu(w_render)),
                Emu(px_emu(h + (4 if single_line else 0))),
            )
            tf = tb.text_frame
            remove_textbox_inset(tf)
            tf.word_wrap = not single_line
            tf.vertical_anchor = MSO_ANCHOR.TOP

            text_content = t.get("text", "")
            if t.get("textTransform") == "uppercase":
                text_content = text_content.upper()
            elif t.get("textTransform") == "lowercase":
                text_content = text_content.lower()

            lines = text_content.split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = align_for(t.get("textAlign", "left"))
                lh = parse_line_height(t.get("lineHeight", ""), t.get("fontSize", 16))
                if lh:
                    p.line_spacing = lh
                run = p.add_run()
                run.text = line
                f = run.font
                f.name = FALLBACK_FONT
                fs = t.get("fontSize") or 16
                f.size = Pt(round(fs * PX_TO_PT, 1))
                f.bold = weight_is_bold(t.get("fontWeight", ""))
                f.italic = (t.get("fontStyle") == "italic")
                col = parse_color(t.get("color") or "")
                if col:
                    # Translucent text colors (e.g. rgba(255,255,255,.72)) blended
                    # over slide canvas so the visible tone survives in PPTX.
                    f.color.rgb = blend_over(col, slide_bg_rgb)
                ls = parse_letter_spacing(t.get("letterSpacing", ""))
                if ls is not None:
                    # python-pptx exposes spacing via XML — set via font._rPr.
                    try:
                        from pptx.oxml.ns import qn
                        rPr = run._r.get_or_add_rPr()
                        # Spacing in PPT is in 1/100 points
                        rPr.set("spc", str(int(round(ls * PX_TO_PT * 100))))
                    except Exception:
                        pass

        # ---- speaker notes (best-effort: page number) --------------------
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.text = f"Slide {slide_n} — exported from Slidev"

    OUT_FILE.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_FILE)
    print(f"Wrote {OUT_FILE} ({OUT_FILE.stat().st_size/1024/1024:.1f} MB)")


if __name__ == "__main__":
    build()
